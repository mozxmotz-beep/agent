from pathlib import Path

from pptx import Presentation

from ppt_agent.models import DeckOutline, SlideSpec
from ppt_agent.renderer import PPTXRenderer, _hex_to_rgb


def test_renderer_output(tmp_path: Path) -> None:
    outline = DeckOutline(
        topic="demo",
        template="modern",
        background="dark",
        slides=[
            SlideSpec(title="封面", subtitle="副标题", layout="title", section="开场"),
            SlideSpec(title="要点", bullets=["A", "B"], layout="title_and_content", section="分析"),
            SlideSpec(title="结论", bullets=["C"], layout="title_and_content", section="结论"),
        ],
    )
    renderer = PPTXRenderer()
    target = tmp_path / "deck.pptx"
    renderer.render(outline, target)
    assert target.exists()

    result = Presentation(str(target))
    assert result.slides[0].shapes.title.text == "封面"
    assert result.slides[1].shapes.title.text.startswith("分析｜")


def test_hex_to_rgb() -> None:
    rgb = _hex_to_rgb("#3B82F6")
    assert (rgb[0], rgb[1], rgb[2]) == (59, 130, 246)
