from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from .models import DeckOutline, SlideSpec


def _hex_to_rgb(color: str) -> RGBColor:
    code = color.lstrip("#")
    return RGBColor(int(code[0:2], 16), int(code[2:4], 16), int(code[4:6], 16))


@dataclass(frozen=True)
class ThemeSpec:
    bg_color: RGBColor
    title_color: RGBColor
    body_color: RGBColor
    accent_color: RGBColor


THEMES: dict[tuple[str, str], ThemeSpec] = {
    ("consulting", "light"): ThemeSpec(
        bg_color=RGBColor(248, 250, 252),
        title_color=RGBColor(15, 23, 42),
        body_color=RGBColor(30, 41, 59),
        accent_color=RGBColor(37, 99, 235),
    ),
    ("consulting", "dark"): ThemeSpec(
        bg_color=RGBColor(15, 23, 42),
        title_color=RGBColor(241, 245, 249),
        body_color=RGBColor(203, 213, 225),
        accent_color=RGBColor(56, 189, 248),
    ),
    ("consulting", "gradient"): ThemeSpec(
        bg_color=RGBColor(239, 246, 255),
        title_color=RGBColor(3, 7, 18),
        body_color=RGBColor(51, 65, 85),
        accent_color=RGBColor(79, 70, 229),
    ),
    ("modern", "light"): ThemeSpec(
        bg_color=RGBColor(255, 255, 255),
        title_color=RGBColor(17, 24, 39),
        body_color=RGBColor(55, 65, 81),
        accent_color=RGBColor(14, 165, 233),
    ),
    ("modern", "dark"): ThemeSpec(
        bg_color=RGBColor(2, 6, 23),
        title_color=RGBColor(226, 232, 240),
        body_color=RGBColor(148, 163, 184),
        accent_color=RGBColor(99, 102, 241),
    ),
    ("modern", "gradient"): ThemeSpec(
        bg_color=RGBColor(236, 253, 245),
        title_color=RGBColor(6, 78, 59),
        body_color=RGBColor(6, 95, 70),
        accent_color=RGBColor(5, 150, 105),
    ),
    ("minimal", "light"): ThemeSpec(
        bg_color=RGBColor(250, 250, 250),
        title_color=RGBColor(38, 38, 38),
        body_color=RGBColor(64, 64, 64),
        accent_color=RGBColor(115, 115, 115),
    ),
    ("minimal", "dark"): ThemeSpec(
        bg_color=RGBColor(23, 23, 23),
        title_color=RGBColor(245, 245, 245),
        body_color=RGBColor(212, 212, 212),
        accent_color=RGBColor(163, 163, 163),
    ),
    ("minimal", "gradient"): ThemeSpec(
        bg_color=RGBColor(245, 245, 244),
        title_color=RGBColor(41, 37, 36),
        body_color=RGBColor(68, 64, 60),
        accent_color=RGBColor(120, 113, 108),
    ),
}


class PPTXRenderer:
    """Render structured slide specs to PPTX."""

    def __init__(self) -> None:
        self.presentation = Presentation()

    def render(self, outline: DeckOutline, output_path: Path) -> Path:
        theme = THEMES[(outline.template, outline.background)]
        if outline.theme_color:
            theme = ThemeSpec(
                bg_color=theme.bg_color,
                title_color=theme.title_color,
                body_color=theme.body_color,
                accent_color=_hex_to_rgb(outline.theme_color),
            )
        for index, slide in enumerate(outline.slides):
            self._render_slide(index=index, slide=slide, theme=theme)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.presentation.save(str(output_path))
        return output_path

    def _render_slide(self, index: int, slide: SlideSpec, theme: ThemeSpec) -> None:
        if index == 0 or slide.layout == "title":
            layout = self.presentation.slide_layouts[0]
            ppt_slide = self.presentation.slides.add_slide(layout)
            self._apply_background(ppt_slide, theme)
            title_shape = ppt_slide.shapes.title
            title_shape.text = slide.title
            self._format_title(title_shape.text_frame, theme)
            if slide.subtitle and len(ppt_slide.placeholders) > 1:
                subtitle_shape = ppt_slide.placeholders[1]
                subtitle = slide.subtitle
                if slide.section:
                    subtitle = f"{slide.section}｜{subtitle}"
                subtitle_shape.text = subtitle
                self._format_subtitle(subtitle_shape.text_frame, theme)
            return

        layout = self.presentation.slide_layouts[1]
        ppt_slide = self.presentation.slides.add_slide(layout)
        self._apply_background(ppt_slide, theme)
        title_text = f"{slide.section}｜{slide.title}" if slide.section else slide.title
        ppt_slide.shapes.title.text = title_text
        self._format_title(ppt_slide.shapes.title.text_frame, theme)

        body = ppt_slide.shapes.placeholders[1].text_frame
        body.clear()
        for line_id, bullet in enumerate(slide.bullets):
            paragraph = body.add_paragraph() if line_id > 0 else body.paragraphs[0]
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = theme.body_color

    def _apply_background(self, ppt_slide, theme: ThemeSpec) -> None:
        fill = ppt_slide.background.fill
        fill.solid()
        fill.fore_color.rgb = theme.bg_color

    def _format_title(self, text_frame, theme: ThemeSpec) -> None:
        for paragraph in text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(34)
            paragraph.font.color.rgb = theme.title_color
            paragraph.alignment = PP_ALIGN.LEFT

    def _format_subtitle(self, text_frame, theme: ThemeSpec) -> None:
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = theme.accent_color
            paragraph.alignment = PP_ALIGN.LEFT
